import streamlit as st
from openai import OpenAI
import json
from pptx import Presentation
import io
import fitz
import base64
import asyncio
from concurrent.futures import ThreadPoolExecutor
import time

# --- Начальная настройка ---
st.set_page_config(
    page_title="Эксперт по подготовке к защите",
    page_icon="🤖",
    layout="wide"
)

EXAMPLE_STORYTELLING_TEXT = """Добрый день, уважаемые коллеги, эксперты, партнёры.
Сегодня мы представляем проект, который находится на пересечении технологий будущего, устойчивого развития и новой философии взаимодействия с мировым океаном.
Проект «Спрут» — это не просто прототип. Это шаг в сторону цивилизованной, экологичной и высокотехнологичной добычи донных полиметаллических конкреций.
Почему это важно?
Полиметаллические конкреции — это настоящие сокровища океанского дна.
Они содержат стратегически важные металлы: марганец, никель, кобальт, медь — именно те, что лежат в основе «зелёной» энергетики, аккумуляторов, микроэлектроники и электротранспорта.
На дне океанов этих ресурсов в 3–4 раза больше, чем на всей суше.
Общие запасы — более 500 миллиардов тонн, из которых половина — полезные минералы.
Это колоссальный потенциал, который может обеспечить человечество сырьём на десятилетия вперёд.
Но есть проблема.
Современные методы добычи устарели и наносят катастрофический урон экосистемам.
При работе тяжёлых драг гибнет более 51% микроорганизмов. Разрушаются биосообщества, нарушается устойчивость экосистем.
При этом спрос на океанические ресурсы не просто растёт — он взрывается.
Количество лицензий на добычу в России с 2020 по 2024 год выросло на 200%.
А мировой рынок подводных технологий демонстрирует рост на 43% в год.
Что мы предложили
Наша задача: создать технологию, которая не разрушает — а бережно взаимодействует с природой.
Изучив существующие решения, проанализировав патенты и собрав межрегиональную команду инженеров, мы разработали принципиально новый подход к подводной добыче.
Так родился «Спрут».
Что такое «Спрут»?
Это мобильная автономная платформа с биомиметическим манипулятором, вдохновлённым природой — щупальцами осьминога и хоботом слона.
Ключевые особенности:
Манипулятор построен по логарифмической спирали — он способен работать с объектами различной формы и диаметра.
Поднимает предметы в 260 раз тяжелее собственного веса.
Управление — как в ручном режиме через приложение, так и в автономном.
Встроенное машинное зрение определяет и классифицирует конкреции прямо на дне.
Что уже сделано
Мы прошли ключевые этапы:
Исследования и патентный анализ
Создание 3D-моделей, разработка электронных компонентов
Сборка и настройка первого прототипа
Первичные испытания в лабораторных условиях
Сегодня у нас — действующий макет манипулятора, рабочая платформа, система управления и программное обеспечение.
Что изменится?
Сравнив нашу систему с традиционными методами, мы получили:
Снижение повреждений донного грунта на 72%
Подъём одной конкреции — за 10 секунд
Скорость передвижения платформы — до 15 км/ч
Автономная работа — до 1,8 часов на одной зарядке
Платформа адаптируется к различным морфологиям дна — мы разработали три её конфигурации.
Кому это выгодно?
Государству — развитие технологического суверенитета и снижение зависимости от импорта
Бизнесу — экологичный имидж, снижение штрафных рисков, экспортный потенциал
Учёным — этичный инструмент для глубоководных исследований
Природе — потому что мы не нарушаем её, а работаем в гармонии с ней
Кто мы?
Мы — команда из 9 инженеров, программистов и схемотехников из разных регионов России.
Нас объединяет страсть к подводной робототехнике и желание переопределить правила в отрасли.
Что дальше?
Проект готов к следующему шагу — испытаниям в реальных морских условиях.
Нам необходимо:
Провести испытания в открытой воде
Дооснастить платформу системой связи и стабилизации
Найти партнёра для запуска пилотного промышленного контракта
Финал
«Спрут» — это не просто машина. Это философия.
Философия бережного, уважительного и умного освоения океана.
Мы верим, что технологии должны быть союзниками природы — а не её врагами.
И если океан — это последнее великое неизведанное пространство на Земле,
мы готовы идти туда. Но идти иначе. С умом, с уважением — и с инновациями.
Спасибо за внимание.
"""

@st.cache_resource
def get_openai_client():
    try:
        client = OpenAI(
            api_key=st.secrets["DEEPSEEK_API_KEY"],
            base_url="https://api.studio.nebius.ai/v1/"
        )
        return client
    except Exception as e:
        st.error(f"Ошибка при инициализации клиента API: {e}")
        return None

client = get_openai_client()

# --- Оптимизированные функции извлечения ---
def extract_text_from_pptx(uploaded_file):
    try:
        pptx_buffer = io.BytesIO(uploaded_file.getvalue())
        prs = Presentation(pptx_buffer)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\n".join(text_runs)
    except Exception as e:
        st.error(f"Ошибка при чтении файла презентации: {e}")
        return ""

def extract_images_from_pptx(uploaded_file):
    images = []
    try:
        pptx_buffer = io.BytesIO(uploaded_file.getvalue())
        prs = Presentation(pptx_buffer)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'shape_type') and shape.shape_type == 13:  # Picture type
                    if hasattr(shape, 'image'):
                        img = shape.image
                        # Ограничиваем размер изображения для ускорения обработки
                        if len(img.blob) < 500000:  # Только изображения меньше 500KB
                            images.append(img.blob)
    except Exception as e:
        st.error(f"Ошибка при извлечении изображений из презентации: {e}")
    return images[:3]  # Ограничиваем до 3 изображений

def extract_images_from_pdf(uploaded_file):
    images = []
    try:
        pdf_bytes = uploaded_file.getvalue()
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            for page_num, page in enumerate(doc):
                if page_num > 5:  # Ограничиваем до 5 страниц
                    break
                for img_ref in page.get_images(full=True):
                    xref = img_ref[0]
                    base_image = doc.extract_image(xref)
                    # Ограничиваем размер изображения
                    if len(base_image["image"]) < 500000:
                        images.append(base_image["image"])
    except Exception as e:
        st.error(f"Ошибка при извлечении изображений из PDF: {e}")
    return images[:3]  # Ограничиваем до 3 изображений

def recognize_images(images: list) -> str:
    descriptions = []
    if not client or not images:
        return ""
    
    # Параллельная обработка изображений
    def process_single_image(img_data, idx):
        try:
            b64 = base64.b64encode(img_data).decode('utf-8')
            # Ограничиваем размер base64 строки
            if len(b64) > 15000:
                b64 = b64[:15000]
                note = " (усечено)"
            else:
                note = ""
            
            prompt = f"Кратко опишите, что изображено на изображении #{idx}{note}."
            
            response = client.chat.completions.create(
                model="google/gemma-3-27b-it",  # Используем указанную модель
                max_tokens=150,  # Увеличено для лучшего качества
                temperature=0.3,
                messages=[{"role": "user", "content": prompt}]
            )
            return f"Изображение #{idx}: {response.choices[0].message.content.strip()}"
        except Exception as e:
            return f"Изображение #{idx}: Ошибка обработки"
    
    # Используем ThreadPoolExecutor для параллельной обработки
    with ThreadPoolExecutor(max_workers=2) as executor:  # Уменьшено до 2 потоков
        futures = [executor.submit(process_single_image, img_bytes, idx) 
                  for idx, img_bytes in enumerate(images[:3], start=1)]
        descriptions = [future.result() for future in futures]
    
    return "\n".join(descriptions)

def extract_text_from_pdf(uploaded_file):
    try:
        pdf_bytes = uploaded_file.getvalue()
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            # Ограничиваем количество страниц для обработки
            pages_to_process = min(10, len(doc))  # Максимум 10 страниц
            return "\n".join([doc[i].get_text() for i in range(pages_to_process)])
    except Exception as e:
        st.error(f"Ошибка при чтении PDF-файла: {e}")
        return ""

def get_analysis_from_deepseek(project_text: str, tone: str):
    if not client:
        return None

    # Ограничиваем длину текста для ускорения обработки
    if len(project_text) > 10000:
        project_text = project_text[:10000] + "\n... (текст усечен для ускорения обработки)"

    storytelling_instruction = f"""
Создайте захватывающий сценарий выступления в стиле {tone.lower()} TED с чётким разделением на три части:
1. "introduction": мощное вступление — обозначьте проблему и представьте название проекта
2. "main_part": развернутая основная часть — опишите цель, ход работы и ключевую новизну проекта
3. "conclusion": убедительное заключение — подведите итоги и подчеркните значимость результатов

(Стиль: живой, конкретный, без абстракций)
"""

    prompt = f"""
Проанализируйте проект и верните строго валидный JSON с:
1. "strengths": 3-5 сильных сторон (массив строк)
2. "weaknesses": 3-5 слабых сторон с рекомендациями (массив строк)
3. "fact_check": 3-4 проверки ключевых утверждений, будь придирчив и въедлив (массив объектов с полями claim, verdict, explanation)
4. "storytelling_script": сценарий выступления ({tone}) - объект с полями introduction, main_part, conclusion
5. "tricky_questions": 4-5 очень каверзных вопросов, как при защите диссертации (массив строк)

Проект:
{project_text}
"""

    try:
        response = client.chat.completions.create(
            model="deepseek-ai/DeepSeek-R1",  # Используем указанную модель
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
            top_p=0.8,
            max_tokens=2000,  # Увеличено для лучшего качества
            response_format={"type": "json_object"}
        )
        return json.loads(response.choices[0].message.content)
    except Exception as e:
        st.error(f"Ошибка при вызове API: {e}")
        return None

# --- Интерфейс приложения ---
st.title("🤖 Эксперт по подготовке к защите")
st.markdown("Загрузите презентацию (`.pdf`, `.pptx`) и/или вставьте текст доклада")

col1, col2 = st.columns(2)

with col1:
    uploaded_file = st.file_uploader("Загрузите презентацию", type=["pptx", "pdf"])

with col2:
    def load_example_text():
        st.session_state.report_text_input = EXAMPLE_STORYTELLING_TEXT

    if "report_text_input" not in st.session_state:
        st.session_state.report_text_input = ""

    report_text = st.text_area(
        "Вставьте текст доклада",
        height=200,
        placeholder="Ваш сценарий выступления...",
        key="report_text_input"
    )

    st.button("✍️ Пример текста", on_click=load_example_text, use_container_width=True)

tone = st.selectbox("🎭 Стиль выступления", ["Вдохновляющий", "Формальный", "Научно-популярный"], index=0)

if st.button("🚀 Проанализировать проект", type="primary", use_container_width=True):
    project_text = ""
    image_descriptions = ""
    
    if uploaded_file is not None:
        with st.spinner("Извлечение данных..."):
            start_time = time.time()
            file_name = uploaded_file.name.lower()
            
            if file_name.endswith('.pptx'):
                images = extract_images_from_pptx(uploaded_file)
                if images:
                    image_descriptions = recognize_images(images)
                project_text = extract_text_from_pptx(uploaded_file)
                
            elif file_name.endswith('.pdf'):
                images = extract_images_from_pdf(uploaded_file)
                if images:
                    image_descriptions = recognize_images(images)
                project_text = extract_text_from_pdf(uploaded_file)
            
            extraction_time = time.time() - start_time
            st.caption(f"⏱️ Извлечение заняло: {extraction_time:.1f} сек")

    # Объединяем данные
    combined_text = "".join(filter(None, [image_descriptions, project_text, report_text]))

    if not combined_text.strip():
        st.warning("Загрузите файл или введите текст")
    else:
        with st.spinner("Анализ ИИ... (~45-90 секунд)"):
            start_time = time.time()
            analysis_result = get_analysis_from_deepseek(combined_text, tone)
            analysis_time = time.time() - start_time
            st.caption(f"⏱️ Анализ занял: {analysis_time:.1f} сек")
            
            if analysis_result:
                # Отображение результатов
                st.success("✅ Анализ завершен!")
                
                tabs = st.tabs(["📊 Сильные/Слабые стороны", "🔍 Фактчек", "🎤 Выступление", "❓ Вопросы"])
                
                with tabs[0]:
                    col1, col2 = st.columns(2)
                    with col1:
                        st.write("**Сильные стороны:**")
                        strengths = analysis_result.get("strengths", [])
                        if isinstance(strengths, list):
                            for strength in strengths:
                                st.success(f"✅ {strength}")
                        else:
                            st.success(f"✅ {strengths}")
                            
                    with col2:
                        st.write("**Для улучшения:**")
                        weaknesses = analysis_result.get("weaknesses", [])
                        if isinstance(weaknesses, list):
                            for weakness in weaknesses:
                                st.warning(f"⚠️ {weakness}")
                        else:
                            st.warning(f"⚠️ {weaknesses}")
                
                with tabs[1]:
                    st.write("**Проверка фактов:**")
                    fact_checks = analysis_result.get("fact_check", [])
                    if isinstance(fact_checks, list):
                        for fact in fact_checks:
                            if isinstance(fact, dict):
                                claim = fact.get('claim', '')
                                verdict = fact.get('verdict', '')
                                explanation = fact.get('explanation', '')
                                with st.expander(f"**{claim}**"):
                                    st.write(f"**Вердикт:** {verdict}")
                                    st.write(f"**Объяснение:** {explanation}")
                            else:
                                st.write(f"**Факт:** {fact}")
                    else:
                        st.write(f"**Факты:** {fact_checks}")
                
                with tabs[2]:
                    st.write("**Сценарий выступления:**")
                    script = analysis_result.get("storytelling_script", {})
                    
                    # Проверяем тип данных script
                    if isinstance(script, str):
                        # Если это строка, пытаемся преобразовать в JSON
                        try:
                            script = json.loads(script)
                        except:
                            st.write(script)
                            script = {}
                    
                    if isinstance(script, dict):
                        sections = [
                            ("🎯 Вступление", "introduction"),
                            ("📝 Основная часть", "main_part"), 
                            ("🏁 Заключение", "conclusion")
                        ]
                        for title, key in sections:
                            content = script.get(key, "")
                            if content:
                                with st.expander(title, expanded=len(content) < 300):
                                    st.write(content)
                    else:
                        st.write("Сценарий:", script)
                
                with tabs[3]:
                    st.write("**Каверзные вопросы:**")
                    questions = analysis_result.get("tricky_questions", [])
                    if isinstance(questions, list):
                        for i, question in enumerate(questions, 1):
                            st.info(f"{i}. {question}")
                    else:
                        st.info(f"Вопросы: {questions}")
            else:
                st.error("Ошибка анализа. Проверьте API ключ и попробуйте позже.")
