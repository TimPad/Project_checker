import streamlit as st
from openai import OpenAI
import json
from pptx import Presentation
import io
import pandas as pd
import fitz  # <<< ИЗМЕНЕНИЕ 1: Импорт новой библиотеки для PDF >>>

# --- Начальная настройка ---
st.set_page_config(
    page_title="Эксперт по подготовке к защите",
    page_icon="🤖",
    layout="wide"
)

# --- Инициализация клиента API ---
@st.cache_resource
def get_openai_client():
    """Инициализирует и возвращает клиент OpenAI для DeepSeek."""
    try:
        client = OpenAI(
            api_key=st.secrets["DEEPSEEK_API_KEY"],
            base_url="https://api.studio.nebius.ai/v1/"
        )
        return client
    except KeyError:
        st.error("Ключ DEEPSEEK_API_KEY не найден в секретах Streamlit. Пожалуйста, добавьте его.")
        return None
    except Exception as e:
        st.error(f"Произошла ошибка при инициализации клиента API: {e}")
        return None

client = get_openai_client()

# --- Вспомогательные функции ---
def extract_text_from_pptx(uploaded_file):
    """Извлекает текст из загруженного файла .pptx."""
    try:
        pptx_buffer = io.BytesIO(uploaded_file.getvalue())
        prs = Presentation(pptx_buffer)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\n".join(text_runs)
    except Exception as e:
        st.error(f"Ошибка при чтении файла презентации: {e}")
        return ""

# <<< ИЗМЕНЕНИЕ 2: Новая функция для извлечения текста из PDF >>>
def extract_text_from_pdf(uploaded_file):
    """Извлекает текст из загруженного файла .pdf."""
    try:
        pdf_bytes = uploaded_file.getvalue()
        # Открываем PDF из байтов в памяти
        with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
            text_pages = []
            for page in doc:
                text_pages.append(page.get_text())
        return "\n".join(text_pages)
    except Exception as e:
        st.error(f"Ошибка при чтении PDF-файла: {e}")
        return ""

def get_analysis_from_deepseek(project_text: str):
    """Отправляет текст в DeepSeek API и получает структурированный анализ."""
    if not client:
        return None

    prompt = f"""
Ты — опытный эксперт по оценке учебно-исследовательских и проектных работ школьников. Твоя роль — быть требовательным, но справедливым наставником. Ты анализируешь предоставленные материалы (текст из презентации и/или доклада) и даешь исчерпывающую, структурированную обратную связь.

Твоя задача — провести комплексный анализ материалов проекта и вернуть результат в виде единого JSON-объекта со следующими ключами:

1.  "strengths": Список (массив строк) сильных сторон проекта.
2.  "weaknesses": Список (массив строк) слабых сторон и зон для улучшения с рекомендациями.
3.  "fact_check": Проверка 3-5 ключевых утверждений. Для каждого утверждения создай объект с ключами "claim", "verdict" ('Истина', 'Ложь' или 'Непроверяемо/Требует уточнения') и "explanation".
4.  "storytelling_script": Объект со сценарием выступления (5-7 мин) с ключами "introduction", "main_part", "conclusion".
5.  "tricky_questions": Список (массив строк) из 5 каверзных вопросов для защиты.

Убедись, что твой ответ — это строго валидный JSON-объект и ничего больше.

Материалы проекта для анализа:
{project_text}
"""

    try:
        response = client.chat.completions.create(
            model="deepseek-ai/DeepSeek-V3",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.6,
            response_format={"type": "json_object"}
        )
        content = response.choices[0].message.content
        return json.loads(content)
    except Exception as e:
        st.error(f"Ошибка при вызове API DeepSeek: {e}")
        return None

# --- Интерфейс приложения Streamlit ---
st.title("🤖 Эксперт по школьным проектам")
# <<< ИЗМЕНЕНИЕ 3: Обновленный текст с упоминанием PDF >>>
st.markdown("Загрузите презентацию (`.pdf`, `.pptx`) и/или вставьте текст доклада, чтобы получить комплексную оценку и рекомендации.")

# Колонки для ввода данных
col1, col2 = st.columns(2)

with col1:
    # <<< ИЗМЕНЕНИЕ 4: Добавлен тип 'pdf' в загрузчик файлов >>>
    uploaded_file = st.file_uploader("Загрузите презентацию (PDF или PPTX)", type=["pptx", "pdf"])

with col2:
    report_text = st.text_area("Вставьте текст доклада", height=250, placeholder="Здесь может быть ваш сценарий выступления, заметки или полный текст доклада...")

# Кнопка для запуска анализа
if st.button("🚀 Провести экспертизу проекта", type="primary", use_container_width=True):
    project_text = ""
    # <<< ИЗМЕНЕНИЕ 5: Логика для определения типа файла и вызова нужной функции >>>
    if uploaded_file is not None:
        with st.spinner("Извлекаю текст из файла..."):
            file_name = uploaded_file.name.lower()
            if file_name.endswith('.pptx'):
                project_text = extract_text_from_pptx(uploaded_file)
            elif file_name.endswith('.pdf'):
                project_text = extract_text_from_pdf(uploaded_file)
    
    combined_text = (project_text + "\n\n" + report_text).strip()

    if not combined_text:
        st.warning("Пожалуйста, загрузите файл или введите текст для анализа.")
    else:
        with st.spinner("ИИ-эксперт изучает ваш проект... Это может занять до 2 минут."):
            analysis_result = get_analysis_from_deepseek(combined_text)

        if analysis_result:
            st.success("Анализ завершен! Вот результаты:")
            
            st.header("📊 Комплексная оценка проекта")

            col_strengths, col_weaknesses = st.columns(2)
            with col_strengths:
                st.subheader("👍 Сильные стороны")
                strengths = analysis_result.get("strengths", [])
                if strengths:
                    for item in strengths:
                        st.markdown(f"- {item}")
                else:
                    st.info("Сильные стороны не определены.")
            
            with col_weaknesses:
                st.subheader("🤔 Слабые стороны и зоны роста")
                weaknesses = analysis_result.get("weaknesses", [])
                if weaknesses:
                    for item in weaknesses:
                        st.markdown(f"- {item}")
                else:
                    st.info("Слабые стороны не определены.")

            st.divider()

            st.header("🔍 Проверка фактов")
            fact_check_data = analysis_result.get("fact_check", [])
            if fact_check_data:
                df = pd.DataFrame(fact_check_data)
                df.columns = ["Утверждение", "Вердикт", "Пояснение"]
                st.dataframe(df, use_container_width=True, hide_index=True)
            else:
                st.info("Не удалось извлечь факты для проверки из текста.")

            st.divider()

            st.header("🎤 Идеальный сценарий выступления (5-7 минут)")
            script = analysis_result.get("storytelling_script")
            
            if isinstance(script, dict):
                with st.expander("Показать/скрыть сценарий", expanded=True):
                    st.subheader("Вступление (≈1 минута)")
                    st.markdown(script.get("introduction", "_Текст для вступления не был сгенерирован._"))
                    
                    st.subheader("Основная часть (≈3-5 минут)")
                    st.markdown(script.get("main_part", "_Текст для основной части не был сгенерирован._"))

                    st.subheader("Заключение (≈1 минута)")
                    st.markdown(script.get("conclusion", "_Текст для заключения не был сгенерирован._"))
            else:
                st.markdown(script or "Не удалось сгенерировать сценарий.")

            st.divider()

            st.header("🤔 Топ-5 каверзных вопросов для защиты")
            questions = analysis_result.get("tricky_questions", [])
            if questions:
                for i, question in enumerate(questions):
                    st.markdown(f"**{i+1}.** {question}")
            else:
                st.info("Не удалось сгенерировать каверзные вопросы.")
                
        else:
            st.error("Не удалось получить результат анализа. Пожалуйста, проверьте консоль на наличие ошибок и попробуйте еще раз.")
